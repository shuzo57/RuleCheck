from io import BytesIO
from pathlib import Path
from typing import IO, Union
from xml.etree import ElementTree as ET

from fastapi import UploadFile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PptxConverter:
    class UnsupportedInput(TypeError):
        pass

    @staticmethod
    def _load_presentation(
        src: Union[str, Path, bytes, bytearray, memoryview, IO[bytes], UploadFile],
    ) -> Presentation:
        if isinstance(src, UploadFile):
            f = src.file
            try: f.seek(0)
            except Exception: pass
            return Presentation(f)
        if isinstance(src, (str, Path)):
            p = Path(src)
            if not p.exists():
                raise FileNotFoundError(p)
            return Presentation(p)
        if isinstance(src, (bytes, bytearray, memoryview)):
            bio = BytesIO(bytes(src))
            bio.seek(0)
            return Presentation(bio)
        if hasattr(src, "read"):
            try:
                src.seek(0)
            except Exception:
                pass
            return Presentation(src)
        raise PptxConverter.UnsupportedInput(type(src))

    @staticmethod
    def convert_to_xml(
        src: Union[str, Path, bytes, bytearray, memoryview, IO[bytes], UploadFile],
        encoding: str = "utf-8",
        pretty: bool = True,
    ) -> str:
        prs = PptxConverter._load_presentation(src)
        root = ET.Element("Document")
        for idx, slide in enumerate(prs.slides, 1):
            s_el = ET.SubElement(root, "Slide", number=str(idx))
            for shape in slide.shapes:
                text = None
                if getattr(shape, "has_text_frame", False):
                    text = (shape.text or "").strip()
                elif hasattr(shape, "text"):
                    text = (shape.text or "").strip()
                if text:
                    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
                    if lines:
                        ET.SubElement(s_el, "Text").text = "\n".join(lines)
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    caption = ""
                    try:
                        alt_nodes = shape._element.xpath(".//p:cNvPr")
                        if alt_nodes:
                            caption = alt_nodes[0].get("descr") or ""
                    except Exception:
                        pass
                    name = getattr(getattr(shape, "image", None), "filename", "image")
                    ET.SubElement(s_el, "Image", name=name, caption=caption)
        if pretty and hasattr(ET, "indent"):
            ET.indent(root, space="  ", level=0)
            return ET.tostring(root, encoding=encoding, xml_declaration=True).decode(encoding)
        if pretty:
            from xml.dom import minidom
            raw = ET.tostring(root, encoding=encoding, xml_declaration=True)
            return minidom.parseString(raw).toprettyxml(indent="  ", encoding=encoding).decode(encoding)
        return ET.tostring(root, encoding=encoding, xml_declaration=True).decode(encoding)
